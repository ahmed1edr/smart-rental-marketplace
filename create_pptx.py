"""
Artisans Connect — PowerPoint Presentation Generator
Creates a professional 15-slide PPTX with a dark premium theme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ─── Color Palette ───
BG_DARK      = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD      = RGBColor(0x1A, 0x1A, 0x35)
BG_TABLE_HDR = RGBColor(0x2D, 0x1B, 0x69)
TEXT_PRIMARY  = RGBColor(0xF0, 0xF0, 0xF5)
TEXT_SECONDARY= RGBColor(0xA0, 0xA0, 0xC0)
TEXT_MUTED    = RGBColor(0x6A, 0x6A, 0x8A)
ACCENT_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
ACCENT_VIOLET = RGBColor(0xA8, 0x55, 0xF7)
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_SUBTLE = RGBColor(0x2A, 0x2A, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set a solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or BG_CARD
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_PRIMARY,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with a single styled paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height):
    """Add an empty text box and return its text_frame for rich text building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, font_size=14, color=TEXT_PRIMARY, bold=False, font_name='Calibri'):
    """Add a styled run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return run


def add_slide_label(slide, text, top=Inches(0.5)):
    """Add the small label at top-left (e.g. '01 — Introduction')."""
    add_text_box(slide, Inches(0.8), top, Inches(4), Inches(0.4),
                 text, font_size=11, color=ACCENT_VIOLET, bold=True, font_name='Calibri')


def add_slide_title(slide, text, top=Inches(0.95)):
    """Add the main slide title."""
    add_text_box(slide, Inches(0.8), top, Inches(11), Inches(0.7),
                 text, font_size=32, color=TEXT_PRIMARY, bold=True, font_name='Calibri')


def add_slide_subtitle(slide, text, top=Inches(1.7)):
    """Add subtitle text."""
    add_text_box(slide, Inches(0.8), top, Inches(9), Inches(0.8),
                 text, font_size=15, color=TEXT_SECONDARY, font_name='Calibri')


def add_accent_line(slide, top=Inches(0.85)):
    """Add a thin accent bar under the label."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(1.8), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_PURPLE
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, items, bullet_color=ACCENT_VIOLET, font_size=13):
    """Add a bullet list to the slide."""
    tf = add_rich_text_box(slide, left, top, width, Inches(len(items) * 0.35 + 0.2))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_run(p, "●  ", font_size=9, color=bullet_color, font_name='Calibri')
        add_run(p, item, font_size=font_size, color=TEXT_SECONDARY, font_name='Calibri')
        p.space_after = Pt(6)
    return tf


def add_card(slide, left, top, width, height, icon, title, description, accent=ACCENT_PURPLE):
    """Add a card-style box with icon, title, and description."""
    card = add_shape(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
    # Accent top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # Icon
    add_text_box(slide, left + Inches(0.2), top + Inches(0.18), Inches(0.6), Inches(0.5),
                 icon, font_size=22, color=accent, alignment=PP_ALIGN.LEFT)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), Inches(0.35),
                 title, font_size=14, color=TEXT_PRIMARY, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(0.92), width - Inches(0.4), height - Inches(1.1),
                 description, font_size=11, color=TEXT_SECONDARY)


def add_table(slide, left, top, width, rows_data, col_widths, header_color=BG_TABLE_HDR):
    """Add a styled table. rows_data[0] is the header row."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(n_rows * 0.38))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.text = cell_text
            p.font.name = 'Calibri'

            if r_idx == 0:
                # Header
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = ACCENT_VIOLET
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_SECONDARY
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if r_idx % 2 == 0 else BG_DARK

            # Vertical center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Padding
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            # Border color
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    ln = tcPr.makeelement(qn(border_name), {})
                    tcPr.append(ln)
                ln.set('w', '6350')
                sf = ln.find(qn('a:solidFill'))
                if sf is None:
                    sf = ln.makeelement(qn('a:solidFill'), {})
                    ln.append(sf)
                srgb = sf.find(qn('a:srgbClr'))
                if srgb is None:
                    srgb = sf.makeelement(qn('a:srgbClr'), {})
                    sf.append(srgb)
                srgb.set('val', '2A2A4A')

    return table_shape


def add_decorative_circle(slide, left, top, size, color, opacity_pct=15):
    """Add a decorative blurred circle for visual interest."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    # Set transparency via XML
    try:
        spPr = shape._element.find(qn('a:solidFill'))
        if spPr is None:
            spPr = shape._element.findall('.//' + qn('a:solidFill'))
            if spPr:
                spPr = spPr[0]
        if spPr is not None:
            srgbClr = spPr.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = srgbClr.makeelement(qn('a:alpha'), {})
                alpha.set('val', str(opacity_pct * 1000))
                srgbClr.append(alpha)
    except Exception:
        pass  # Fallback: just skip transparency
    shape.line.fill.background()
    return shape


def add_flow_step(slide, left, top, number, text, accent=ACCENT_VIOLET):
    """Add a flow step box."""
    box = add_shape(slide, left, top, Inches(1.9), Inches(0.85), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
    add_text_box(slide, left + Inches(0.1), top + Inches(0.08), Inches(1.7), Inches(0.35),
                 str(number), font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_text_box(slide, left + Inches(0.1), top + Inches(0.42), Inches(1.7), Inches(0.4),
                 text, font_size=9, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)


def add_flow_arrow(slide, left, top):
    """Add a → arrow between flow steps."""
    add_text_box(slide, left, top, Inches(0.35), Inches(0.85),
                 "→", font_size=18, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER, font_name='Calibri')


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(-2), Inches(-2), Inches(7), ACCENT_PURPLE, 10)
add_decorative_circle(slide, Inches(9), Inches(4), Inches(6), ACCENT_CYAN, 8)

add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(0.8),
             "🔧", font_size=48, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.4), SLIDE_W, Inches(1),
             "Artisans Connect", font_size=52, color=ACCENT_VIOLET, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Calibri')
add_text_box(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.5),
             "La plateforme qui connecte citoyens et artisans", font_size=20,
             color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.0), SLIDE_W, Inches(0.4),
             "Rapport d'Analyse de Projet", font_size=15,
             color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Authors box
box = add_shape(slide, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.5),
                fill_color=BG_CARD, border_color=ACCENT_PURPLE, border_width=Pt(1))
add_text_box(slide, Inches(4.5), Inches(5.05), Inches(4.3), Inches(0.4),
             "Solayman Elqasmi  •  Hrimech Abdessalam", font_size=12,
             color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.7), SLIDE_W, Inches(0.4),
             "Année académique 2025 – 2026", font_size=13,
             color=ACCENT_VIOLET, alignment=PP_ALIGN.CENTER, font_name='Calibri')


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(10), Inches(-2), Inches(5), ACCENT_PURPLE, 8)

add_slide_label(slide, "01 — INTRODUCTION")
add_accent_line(slide)
add_slide_title(slide, "Qu'est-ce que Artisans Connect ?")
add_slide_subtitle(slide, "Une plateforme numérique de type marketplace conçue pour combler l'écart entre les citoyens cherchant des services artisanaux qualifiés et les artisans offrant ces services au Maroc.")

add_card(slide, Inches(0.8), Inches(3.1), Inches(3.6), Inches(2.2),
         "🔍", "Recherche intelligente",
         "Rechercher des artisans par type de service et ville en quelques clics",
         ACCENT_PURPLE)
add_card(slide, Inches(4.8), Inches(3.1), Inches(3.6), Inches(2.2),
         "⭐", "Profils vérifiés",
         "Consulter les profils, évaluations et avis des artisans de confiance",
         ACCENT_CYAN)
add_card(slide, Inches(8.8), Inches(3.1), Inches(3.6), Inches(2.2),
         "📞", "Contact direct",
         "Entrer en contact direct avec les artisans qualifiés près de chez vous",
         ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Problématique
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(-1), Inches(3), Inches(5), ACCENT_RED, 6)

add_slide_label(slide, "02 — PROBLÉMATIQUE")
add_accent_line(slide)
add_slide_title(slide, "Les défis actuels")
add_slide_subtitle(slide, "Les méthodes traditionnelles — bouche-à-oreille, annonces classées — sont inefficaces et ne garantissent ni la qualité ni la transparence.")

# Citoyen column
add_shape(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(3.8), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(3.05), Inches(5), Inches(0.35),
             "⚠️  Point de vue du Citoyen", font_size=15, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(3.5), Inches(5), [
    "Difficulté à trouver un artisan de confiance à proximité",
    "Manque de transparence des prix",
    "Délais et qualité de service médiocre",
    "Crainte de fraude et d'escroquerie",
    "Absence de système d'évaluation ou de garantie",
], bullet_color=ACCENT_RED)

# Artisan column
add_shape(slide, Inches(6.8), Inches(2.9), Inches(5.5), Inches(3.8), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(7.1), Inches(3.05), Inches(5), Inches(0.35),
             "🔨  Point de vue de l'Artisan", font_size=15, color=ACCENT_AMBER, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(3.5), Inches(5), [
    "Difficulté à trouver de nouveaux clients",
    "Visibilité numérique limitée",
    "Absence de plateforme pour présenter ses travaux",
    "Dépendance aux réseaux informels",
], bullet_color=ACCENT_AMBER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Objectifs
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(9), Inches(-1), Inches(5), ACCENT_GREEN, 8)

add_slide_label(slide, "03 — OBJECTIFS")
add_accent_line(slide)
add_slide_title(slide, "Objectifs du projet")

objectives = [
    ("🔗", "Plateforme de connexion", "Fournir une plateforme simple, sécurisée et accessible reliant citoyens et artisans", ACCENT_PURPLE),
    ("🔍", "Recherche ciblée", "Permettre la recherche d'artisans par type de service et par ville", ACCENT_CYAN),
    ("✅", "Profils vérifiés", "Créer des profils vérifiés avec affichage des travaux antérieurs", ACCENT_GREEN),
    ("⭐", "Système d'évaluation", "Notes sur 5 étoiles et avis pour établir la confiance", ACCENT_AMBER),
    ("💎", "Abonnement Premium", "Modèle d'abonnement pour augmenter la visibilité des artisans", ACCENT_VIOLET),
    ("🌍", "Scalabilité", "Croissance à l'échelle nationale et internationale", ACCENT_CYAN),
]
for i, (icon, title, desc, accent) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.0) + row * Inches(2.5)
    add_card(slide, x, y, Inches(3.7), Inches(2.1), icon, title, desc, accent)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Public cible
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(-1), Inches(-1), Inches(5), ACCENT_CYAN, 7)

add_slide_label(slide, "04 — PUBLIC CIBLE")
add_accent_line(slide)
add_slide_title(slide, "Analyse du public cible")

# Persona card
add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(2.25), Inches(2), Inches(0.5),
             "👤", font_size=36)
add_text_box(slide, Inches(1.1), Inches(2.9), Inches(4.5), Inches(0.4),
             "Persona — Citoyen", font_size=18, color=TEXT_PRIMARY, bold=True)

persona_data = [
    ["Attribut", "Détails"],
    ["Nom", "Abdelrahim (profil représentatif)"],
    ["Âge", "20 – 60 ans"],
    ["Ville", "Safi"],
    ["CSP", "Ouvrier artisanal / Grand public"],
    ["Pouvoir d'achat", "~3 000 DH/mois"],
]
add_table(slide, Inches(1.1), Inches(3.5), Inches(4.6), persona_data,
          [Inches(1.6), Inches(3.0)])

# Attentes
add_shape(slide, Inches(6.5), Inches(2.1), Inches(5.8), Inches(2.0), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(6.8), Inches(2.2), Inches(5), Inches(0.3),
             "✅  Attentes", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(2.55), Inches(5.2), [
    "Trouver rapidement un artisan qualifié",
    "Plateforme simple et facile à utiliser",
    "Tarifs transparents et services fiables",
], bullet_color=ACCENT_GREEN, font_size=11)

# Canaux
add_shape(slide, Inches(6.5), Inches(4.4), Inches(5.8), Inches(2.2), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(6.8), Inches(4.5), Inches(5), Inches(0.3),
             "📢  Canaux d'acquisition", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(4.85), Inches(5.2), [
    "Réseaux sociaux (Facebook, Instagram)",
    "Groupes et messages WhatsApp",
    "Campagnes publicitaires en ligne",
    "Hub Hihof (centre numérique local)",
], bullet_color=ACCENT_CYAN, font_size=11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Exigences fonctionnelles
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(10), Inches(4), Inches(5), ACCENT_AMBER, 7)

add_slide_label(slide, "05 — EXIGENCES")
add_accent_line(slide)
add_slide_title(slide, "Exigences fonctionnelles")

# Citoyens
add_card(slide, Inches(0.8), Inches(2.1), Inches(3.6), Inches(4.6),
         "👨‍👩‍👧‍👦", "Pour les Citoyens", "", ACCENT_CYAN)
add_bullet_list(slide, Inches(1.0), Inches(3.5), Inches(3.2), [
    "S'inscrire et créer un compte",
    "Rechercher par service et ville",
    "Consulter les profils détaillés",
    "Comparer services et prix",
    "Contacter un artisan directement",
    "Laisser des évaluations et notes",
], bullet_color=ACCENT_CYAN, font_size=11)

# Artisans
add_card(slide, Inches(4.8), Inches(2.1), Inches(3.6), Inches(4.6),
         "🔨", "Pour les Artisans", "", ACCENT_AMBER)
add_bullet_list(slide, Inches(5.0), Inches(3.5), Inches(3.2), [
    "Créer un compte artisan",
    "Compléter le profil détaillé",
    "Télécharger photos de travaux",
    "Accès abonnement Premium",
    "Profil visible une fois complété",
], bullet_color=ACCENT_AMBER, font_size=11)

# Admins
add_card(slide, Inches(8.8), Inches(2.1), Inches(3.6), Inches(4.6),
         "⚙️", "Pour les Admins", "", ACCENT_RED)
add_bullet_list(slide, Inches(9.0), Inches(3.5), Inches(3.2), [
    "Tableau de bord de gestion",
    "Modération profils et avis",
    "Statistiques et analyses",
    "Gestion des abonnements",
], bullet_color=ACCENT_RED, font_size=11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Structure du site
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(-2), Inches(4), Inches(5), ACCENT_VIOLET, 7)

add_slide_label(slide, "06 — STRUCTURE")
add_accent_line(slide)
add_slide_title(slide, "Structure du site web")
add_slide_subtitle(slide, "Les pages principales et les flux utilisateur de la plateforme")

pages = [
    ("🏠", "Page d'accueil", "Recherche, artisans vedettes, catégories", ACCENT_PURPLE),
    ("📝", "Inscription", "Créer un compte citoyen ou artisan", ACCENT_CYAN),
    ("🔐", "Connexion", "Authentification email / mot de passe", ACCENT_GREEN),
    ("👷", "Profil Artisan", "Portfolio, évaluations, contact", ACCENT_AMBER),
    ("📋", "Résultats", "Liste filtrée des artisans", ACCENT_RED),
    ("📊", "Dashboard Admin", "Interface de gestion backend", ACCENT_VIOLET),
]
for i, (icon, title, desc, accent) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.8) + row * Inches(2.2)
    add_card(slide, x, y, Inches(3.7), Inches(1.8), icon, title, desc, accent)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Parcours utilisateur
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(10), Inches(-1), Inches(5), ACCENT_CYAN, 7)

add_slide_label(slide, "07 — PARCOURS UTILISATEUR")
add_accent_line(slide)
add_slide_title(slide, "Parcours utilisateur")

# Citoyen flow
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5), Inches(0.35),
             "🔵  Parcours Citoyen", font_size=15, color=ACCENT_CYAN, bold=True)

citoyen_steps = ["Page d'accueil", "Recherche", "Résultats", "Sélection", "Contact"]
for i, step_text in enumerate(citoyen_steps):
    x = Inches(0.8) + i * Inches(2.25)
    add_flow_step(slide, x, Inches(2.6), i + 1, step_text, ACCENT_CYAN)
    if i < 4:
        add_flow_arrow(slide, x + Inches(1.9), Inches(2.6))

# Artisan flow
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.35),
             "🟠  Parcours Artisan", font_size=15, color=ACCENT_AMBER, bold=True)

artisan_steps = ["Crée compte", "Se connecte", "Complète profil", "Upload photos", "Profil visible ✓"]
for i, step_text in enumerate(artisan_steps):
    x = Inches(0.8) + i * Inches(2.25)
    add_flow_step(slide, x, Inches(4.4), i + 1, step_text, ACCENT_AMBER)
    if i < 4:
        add_flow_arrow(slide, x + Inches(1.9), Inches(4.4))

# USP box
add_shape(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(1.0),
          fill_color=BG_CARD, border_color=ACCENT_PURPLE, border_width=Pt(1))
add_text_box(slide, Inches(1.8), Inches(5.75), Inches(9.7), Inches(0.5),
             "« Une plateforme simple et sécurisée qui permet aux citoyens de trouver rapidement des artisans fiables, avec des prix transparents et des avis vérifiés. »",
             font_size=12, color=ACCENT_VIOLET, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.8), Inches(6.25), Inches(9.7), Inches(0.3),
             "— Proposition de Valeur Unique (USP)", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Fonctionnalités
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_slide_label(slide, "08 — FONCTIONNALITÉS")
add_accent_line(slide)
add_slide_title(slide, "Fonctionnalités principales")

features_data = [
    ["Fonctionnalité", "Description", "Rôle"],
    ["🔍 Recherche avancée", "Rechercher par type de service et ville", "Citoyen"],
    ["👷 Profil artisan", "Profil détaillé avec photos, notes, coordonnées", "Citoyen / Artisan"],
    ["📝 Inscription & connexion", "Création de compte avec sélection du rôle", "Tous"],
    ["⭐ Notes et avis", "Publication de notes sur 5 étoiles et commentaires", "Citoyen"],
    ["💎 Abonnement premium", "Plan payant pour augmenter la visibilité", "Artisan"],
    ["📊 Dashboard admin", "Gestion des utilisateurs, contenu et statistiques", "Administrateur"],
    ["📞 Contact direct", "Affichage du numéro pour communication directe", "Citoyen"],
    ["📸 Portfolio", "Télécharger et afficher des photos de projets", "Artisan"],
]
add_table(slide, Inches(0.8), Inches(2.1), Inches(11.5), features_data,
          [Inches(3.2), Inches(6.0), Inches(2.3)])


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Technologies
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(9), Inches(3), Inches(5), ACCENT_GREEN, 7)

add_slide_label(slide, "09 — TECHNOLOGIES")
add_accent_line(slide)
add_slide_title(slide, "Stack technique")

tech_data = [
    ["Couche", "Technologie", "Objectif"],
    ["Frontend", "HTML, CSS, JavaScript (React / Vue)", "Interface utilisateur et interactions"],
    ["Backend", "Node.js / PHP", "Logique serveur, gestion des API"],
    ["Base de données", "MySQL / MongoDB", "Stockage et gestion des données"],
    ["Hébergement", "VPS haute performance", "Hébergement serveur professionnel"],
    ["Sécurité", "SSL/TLS, Chiffrement", "Protection des données"],
    ["Domaine", ".com ou .ma", "Adresse Web"],
]
add_table(slide, Inches(0.8), Inches(2.1), Inches(7), tech_data,
          [Inches(1.8), Inches(3.0), Inches(2.2)])

# Security section
add_shape(slide, Inches(8.3), Inches(2.1), Inches(4.3), Inches(4.2), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
add_text_box(slide, Inches(8.6), Inches(2.25), Inches(3.8), Inches(0.35),
             "🔒  Mesures de sécurité", font_size=14, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(8.6), Inches(2.7), Inches(3.8), [
    "Protection : SQL Injection, XSS, CSRF",
    "Chiffrement des mots de passe",
    "Connexion HTTPS sécurisée",
    "Sauvegardes automatisées",
], bullet_color=ACCENT_RED, font_size=11)

add_text_box(slide, Inches(8.6), Inches(4.3), Inches(3.8), Inches(0.35),
             "🛠️  Outils de développement", font_size=14, color=ACCENT_CYAN, bold=True)
add_bullet_list(slide, Inches(8.6), Inches(4.7), Inches(3.8), [
    "GitHub Pro",
    "Suites de test automatisées",
    "Plugins de développement premium",
], bullet_color=ACCENT_CYAN, font_size=11)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Base de données
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(-1), Inches(4), Inches(5), ACCENT_AMBER, 6)

add_slide_label(slide, "10 — BASE DE DONNÉES")
add_accent_line(slide)
add_slide_title(slide, "Conception de la base de données")

entities = [
    ("👤", "Utilisateur (Citoyen)", "ID, Prénom, Nom, Email, Mot de passe, Téléphone, Ville, Âge, Type de compte", ACCENT_CYAN),
    ("🔨", "Artisan", "ID, Prénom, Nom, Métier, Ville, Téléphone, Photo profil, Photos travaux, Prix, Expérience, Évaluation", ACCENT_AMBER),
    ("⭐", "Avis", "ID, ID Citoyen, ID Artisan, Note (1-5), Commentaire, Date", ACCENT_GREEN),
    ("💎", "Abonnement", "ID, ID Artisan, Type plan, Date début, Date fin, Statut paiement", ACCENT_VIOLET),
    ("⚙️", "Administrateur", "ID, Nom, Email, Mot de passe, Rôle", ACCENT_RED),
]

for i, (icon, name, fields, accent) in enumerate(entities):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(2.1) + row * Inches(1.7)
    w = Inches(5.8)
    h = Inches(1.4)

    add_shape(slide, x, y, w, h, fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = accent
    top_line.line.fill.background()

    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.35),
                 f"{icon}  {name}", font_size=14, color=accent, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.8),
                 fields, font_size=10, color=TEXT_SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Analyse concurrentielle
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(10), Inches(-1), Inches(5), ACCENT_GREEN, 8)

add_slide_label(slide, "11 — ANALYSE CONCURRENTIELLE")
add_accent_line(slide)
add_slide_title(slide, "Positionnement concurrentiel")

comp_data = [
    ["Critères", "HirafiMall (Existant)", "Artisans Connect ✨"],
    ["🌍 Portée", "Distribution limitée", "Échelle nationale & internationale"],
    ["🎨 Interface", "Simple mais peu attrayante", "Moderne & professionnelle"],
    ["💎 Abonnement", "Aucun Premium", "Premium disponible"],
    ["🖌️ Design", "Design basique", "Design professionnel haut de gamme"],
    ["📈 Scalabilité", "Croissance limitée", "Architecture scalable"],
]
add_table(slide, Inches(0.8), Inches(2.2), Inches(11.5), comp_data,
          [Inches(2.5), Inches(4.5), Inches(4.5)])

# Conclusion box
add_shape(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.8),
          fill_color=BG_CARD, border_color=ACCENT_GREEN, border_width=Pt(1))
add_text_box(slide, Inches(1.8), Inches(5.35), Inches(9.7), Inches(0.7),
             "✅  Artisans Connect se positionne comme une alternative supérieure combinant facilité d'utilisation, design moderne, scalabilité géographique et monétisation par abonnement.",
             font_size=12, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — SWOT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_slide_label(slide, "12 — ANALYSE SWOT")
add_accent_line(slide)
add_slide_title(slide, "Analyse SWOT")

swot = [
    ("💪 Forces", ACCENT_GREEN, [
        "Interface simple et conviviale",
        "Profils d'artisans vérifiés",
        "Tarification transparente",
        "Système d'évaluation et d'avis",
        "Monétisation par abonnement premium",
    ]),
    ("⚠️ Faiblesses", ACCENT_RED, [
        "Plateforme nouvelle sans réputation",
        "Investissement marketing initial significatif",
        "Dépendance du taux d'adoption des artisans",
    ]),
    ("🚀 Opportunités", ACCENT_CYAN, [
        "Demande croissante de plateformes en ligne",
        "Large marché informel d'artisans au Maroc",
        "Potentiel d'expansion nationale / internationale",
        "Concurrence numérique directe limitée",
    ]),
    ("⚡ Menaces", ACCENT_AMBER, [
        "Préoccupations de confiance des utilisateurs",
        "Concurrence des acteurs établis",
        "Barrières à l'adoption numérique",
        "Facteurs économiques / pouvoir d'achat",
    ]),
]

positions = [(Inches(0.8), Inches(2.1)), (Inches(6.6), Inches(2.1)),
             (Inches(0.8), Inches(4.6)), (Inches(6.6), Inches(4.6))]

for (title, color, items), (x, y) in zip(swot, positions):
    w, h = Inches(5.5), Inches(2.2)
    add_shape(slide, x, y, w, h, fill_color=BG_CARD, border_color=color, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.45), w - Inches(0.4),
                    items, bullet_color=color, font_size=10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Budget
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(10), Inches(4), Inches(5), ACCENT_AMBER, 7)

add_slide_label(slide, "13 — BUDGET")
add_accent_line(slide)
add_slide_title(slide, "Budget & Ressources")

# Technical resources
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(3), Inches(0.3),
             "💻 Ressources Techniques", font_size=12, color=ACCENT_CYAN, bold=True)
tech_budget = [
    ["Élément", "Coût (DH)"],
    ["Développement Full-stack", "15 000 – 20 000"],
    ["Hébergement VPS (annuel)", "2 500 – 4 000"],
    ["Nom de domaine", "150 – 200"],
    ["Base de données", "0 – 500"],
    ["Cybersécurité", "500 – 1 000"],
    ["Outils développement", "600 – 1 500"],
]
add_table(slide, Inches(0.8), Inches(2.35), Inches(5.5), tech_budget, [Inches(3.0), Inches(2.5)])

# Human resources
add_text_box(slide, Inches(6.8), Inches(2.0), Inches(3), Inches(0.3),
             "👥 Ressources Humaines (mensuel)", font_size=12, color=ACCENT_VIOLET, bold=True)
hr_budget = [
    ["Rôle", "Coût (DH)"],
    ["Développeur Web", "10 000"],
    ["Designer UX/UI", "3 000 – 5 000"],
    ["Gestionnaire de contenu", "3 000 – 4 800"],
    ["Support technique", "2 500 – 4 000"],
]
add_table(slide, Inches(6.8), Inches(2.35), Inches(5.5), hr_budget, [Inches(3.0), Inches(2.5)])

# Material + Intangible
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(3), Inches(0.3),
             "🖥️ Ressources Matérielles", font_size=12, color=ACCENT_AMBER, bold=True)
mat_budget = [
    ["Élément", "Coût (DH)"],
    ["Ordinateur professionnel", "8 000 – 12 000"],
    ["Internet (mensuel)", "300 – 500"],
    ["Mobile professionnel", "1 500 – 3 000"],
]
add_table(slide, Inches(0.8), Inches(5.65), Inches(5.5), mat_budget, [Inches(3.0), Inches(2.5)])

add_text_box(slide, Inches(6.8), Inches(5.3), Inches(3), Inches(0.3),
             "✨ Ressources Intangibles", font_size=12, color=ACCENT_GREEN, bold=True)
intang_budget = [
    ["Élément", "Coût (DH)"],
    ["Identité de marque", "500"],
    ["Contenu du site", "1 000"],
    ["Marketing de lancement", "6 000"],
]
add_table(slide, Inches(6.8), Inches(5.65), Inches(5.5), intang_budget, [Inches(3.0), Inches(2.5)])


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_decorative_circle(slide, Inches(4), Inches(0), Inches(6), ACCENT_PURPLE, 8)
add_decorative_circle(slide, Inches(-1), Inches(4), Inches(5), ACCENT_CYAN, 6)
add_decorative_circle(slide, Inches(10), Inches(3), Inches(4), ACCENT_GREEN, 6)

add_slide_label(slide, "14 — CONCLUSION")
add_accent_line(slide)

add_text_box(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.6),
             "Artisans Connect — Prêt pour le lancement", font_size=32,
             color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')

# Stats boxes
stats = [
    ("3", "Rôles utilisateur", ACCENT_VIOLET),
    ("8+", "Fonctionnalités clés", ACCENT_CYAN),
    ("5", "Entités de données", ACCENT_GREEN),
    ("~50K DH", "Budget estimé", ACCENT_AMBER),
]
for i, (val, label, color) in enumerate(stats):
    x = Inches(1.0) + i * Inches(2.95)
    add_shape(slide, x, Inches(2.0), Inches(2.5), Inches(1.2), fill_color=BG_CARD, border_color=BORDER_SUBTLE, border_width=Pt(1))
    add_text_box(slide, x, Inches(2.1), Inches(2.5), Inches(0.6),
                 val, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name='Calibri')
    add_text_box(slide, x, Inches(2.65), Inches(2.5), Inches(0.3),
                 label, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Conclusion box
add_shape(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(2.0),
          fill_color=BG_CARD, border_color=ACCENT_PURPLE, border_width=Pt(1))

add_text_box(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(0.5),
             "✅", font_size=36, alignment=PP_ALIGN.CENTER)

tf = add_rich_text_box(slide, Inches(2), Inches(4.25), Inches(9.3), Inches(0.8))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_run(p, "Artisans Connect est une initiative numérique bien fondée qui répond à un besoin réel du marché marocain. ", font_size=12, color=TEXT_SECONDARY)
add_run(p, "Un écosystème équilibré et durable.", font_size=12, color=ACCENT_VIOLET, bold=True)

add_text_box(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.3),
             "Ce rapport confirme la viabilité comme initiative socialement impactante et commercialement durable.",
             font_size=11, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# Merci
add_text_box(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.6),
             "Merci ! 🙏", font_size=28, color=TEXT_MUTED, bold=True,
             alignment=PP_ALIGN.CENTER, font_name='Calibri')
add_text_box(slide, Inches(0), Inches(6.35), SLIDE_W, Inches(0.35),
             "Solayman Elqasmi  •  Hrimech Abdessalam  •  2025-2026",
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = "Artisans_Connect_Presentation.pptx"
prs.save(output_path)
print(f"[OK] Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
